# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
YELLOW1 = RGBColor(0xFF, 0xD7, 0x00)  # Gold
YELLOW2 = RGBColor(0xFF, 0xF3, 0xCD)  # LemonChiffon
YELLOW3 = RGBColor(0xFF, 0xA5, 0x00)  # Orange
BLUE1 = RGBColor(0x1E, 0x90, 0xFF)    # DodgerBlue
BLUE2 = RGBColor(0x41, 0x69, 0xE1)   # RoyalBlue
BLUE3 = RGBColor(0x87, 0xCE, 0xEB)   # LightSkyBlue
DARK = RGBColor(0x2C, 0x3E, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xFA, 0xF8, 0xF0)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
LOGO_PATH = None  # Add logo path if available
LOSS_CURVE = os.path.join(BASE_DIR, "lora_output", "loss_curve.png")
EVAL_GRAPH = os.path.join(BASE_DIR, "lora_output", "eval_graph.png")

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name="Microsoft JhengHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_slide(slide, items, left=Inches(1), top=Inches(2.5), width=Inches(11), font_size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft JhengHei"
        p.space_after = Pt(8)
    return txBox


# ============ Slide 1: Title ============
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, WHITE)

# Top yellow bar
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.15), YELLOW1)
# Bottom blue bar
add_shape(slide, Inches(0), Inches(7.35), Inches(13.33), Inches(0.15), BLUE2)

# Large blue block on left
add_shape(slide, Inches(0.5), Inches(1.5), Inches(0.15), Inches(4.5), BLUE2, MSO_SHAPE.RECTANGLE)

add_text_box(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(1.5),
             "OpenClaw FAQ 智能問答系統", font_size=44, bold=True, color=BLUE2)

add_text_box(slide, Inches(1.2), Inches(3.5), Inches(10), Inches(1),
             "基於 Qwen3.5-0.8B-Base 的 LoRA 微調實踐", font_size=28, color=BLUE1)

add_text_box(slide, Inches(1.2), Inches(5.0), Inches(10), Inches(1),
             "NLP 微調項目報告  |  2026", font_size=20, color=DARK)

# Yellow accent box
add_shape(slide, Inches(1.2), Inches(6.2), Inches(4), Inches(0.06), YELLOW1)


# ============ Slide 2: Project Overview ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), YELLOW1)
add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), BLUE2, MSO_SHAPE.RECTANGLE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "項目概述", font_size=36, bold=True, color=BLUE2)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), YELLOW1)

items = [
    "目標：對 Qwen3.5-0.8B-Base 進行指令微調，使其掌握 OpenClaw 產品的 FAQ 知識",
    "資料：208 組人工整理的 OpenClaw 問答對（中英文混合）",
    "方法：使用 LoRA（Low-Rank Adaptation）進行參數高效微調",
    "基礎模型：Qwen3.5-0.8B-Base（8.6 億參數，支援多模態）",
    "訓練資源：單張 NVIDIA L20（48GB）",
    "服務化：提供 OpenAI 兼容 API，支援即時問答",
]
add_bullet_slide(slide, items, font_size=18)

# Info cards
info_data = [("模型大小", "0.8B"), ("訓練資料", "207筆"), ("LoRA Rank", "8"), ("訓練輪次", "30")]
for i, (label, val) in enumerate(info_data):
    x = Inches(0.8 + i * 3.1)
    card = add_shape(slide, x, Inches(5.2), Inches(2.8), Inches(1.5), BLUE3)
    add_text_box(slide, x + Inches(0.2), Inches(5.3), Inches(2.4), Inches(0.4),
                 label, font_size=14, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(5.7), Inches(2.4), Inches(0.8),
                 val, font_size=28, bold=True, color=BLUE2, align=PP_ALIGN.CENTER)


# ============ Slide 3: Model ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), YELLOW1)
add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), BLUE2, MSO_SHAPE.RECTANGLE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "基礎模型：Qwen3.5-0.8B-Base", font_size=36, bold=True, color=BLUE2)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), YELLOW1)

items = [
    "架構：混合注意力機制（Hybrid Attention）— 18層 Linear Attention + 6層 Full Attention",
    "上下文長度：最高 262,144 tokens（超長上下文支援）",
    "詞表大小：248,320 tokens，支援中英文混合輸入",
    "多模態：支援文字 + 圖片 + 影片輸入（本項目僅使用文字）",
    "MTP（Multi-Token Prediction）：支援多 token 同時預測，提升推理效率",
]
add_bullet_slide(slide, items, font_size=18)

# Architecture diagram
arch_items = [
    ("24層 Decoder", "6層 Full Attention\n18層 Linear Attention"),
    ("嵌入維度", "1024"),
    ("注意力頭", "8 heads, KV 2 heads"),
    ("MLP", "Gate/SiLU 門控結構"),
]
for i, (k, v) in enumerate(arch_items):
    y = Inches(4.8 + i * 0.65)
    add_shape(slide, Inches(0.8), y, Inches(3), Inches(0.5), YELLOW2)
    add_text_box(slide, Inches(1), y, Inches(2.6), Inches(0.5), f"{k}：{v}", font_size=14, color=DARK)


# ============ Slide 4: Dataset ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), YELLOW1)
add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), BLUE2, MSO_SHAPE.RECTANGLE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "資料集說明", font_size=36, bold=True, color=BLUE2)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), YELLOW1)

items = [
    "來源：OpenClaw 官方 FAQ 文檔及社群常見問題",
    "總量：208 條 QA 對（過濾後 207 條有效）",
    "訓練集：175 條（85%）   驗證集：32 條（15%）",
    "覆蓋範圍：安裝部署、模型配置、管道集成、記憶系統、安全性、故障排除等",
    "資料格式：問題（query）+ 答案（pos），中英文混合",
    "範例：\n     Q：OpenClaw 要錢嗎？\n     A：OpenClaw 軟體本身完全免費且開源...",
]
add_bullet_slide(slide, items, font_size=17, top=Inches(1.6))


# ============ Slide 5: LoRA Method ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), YELLOW1)
add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), BLUE2, MSO_SHAPE.RECTANGLE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "LoRA 微調方法", font_size=36, bold=True, color=BLUE2)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), YELLOW1)

items = [
    "LoRA（Low-Rank Adaptation）— 在原始權重旁插入低秩矩陣，僅訓練新增參數",
    "目標模組：q_proj, k_proj, v_proj, o_proj（注意力層）+ gate/up/down_proj（MLP層）",
    "LoRA 配置： rank=8, alpha=16, dropout=0.05",
    "只訓練語言模型部分，凍結視覺編碼器（Vision Encoder）",
    "訓練參數：約 319 萬（僅佔總參數的 0.37%）",
    "優化器：AdamW（lr=2e-4, weight_decay=0.01）",
    "排程：線性 warmup（10%）+ 線性衰減",
]
add_bullet_slide(slide, items, font_size=17)

# Highlight box
add_shape(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1), BLUE3)
add_text_box(slide, Inches(1), Inches(5.9), Inches(10.6), Inches(0.8),
             "W = W₀ + ΔW         ΔW = B × A         B ∈ ℝ^(d×r), A ∈ ℝ^(r×k), r ≪ min(d, k)",
             font_size=18, bold=True, color=BLUE2, align=PP_ALIGN.CENTER)


# ============ Slide 6: Loss Curve ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), YELLOW1)
add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), BLUE2, MSO_SHAPE.RECTANGLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
             "訓練損失曲線", font_size=36, bold=True, color=BLUE2)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), YELLOW1)

if os.path.exists(LOSS_CURVE):
    slide.shapes.add_picture(LOSS_CURVE, Inches(0.8), Inches(1.3), Inches(7), Inches(5.5))

# Observations
obs = [
    "Train Loss 持續下降至 ~0.013",
    "Val Loss 第3輪達最低 0.76",
    "之後 Val Loss 反彈 → 過擬合",
]
for i, o in enumerate(obs):
    y = Inches(1.5 + i * 0.7)
    add_shape(slide, Inches(8.5), y, Inches(4.2), Inches(0.55), YELLOW2 if i < 2 else RGBColor(0xFF, 0xCC, 0xCB))
    add_text_box(slide, Inches(8.7), y + Inches(0.05), Inches(3.8), Inches(0.5),
                 o, font_size=15, color=DARK)


# ============ Slide 7: Evaluation ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), YELLOW1)
add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), BLUE2, MSO_SHAPE.RECTANGLE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
             "評估與問題分析", font_size=36, bold=True, color=BLUE2)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), YELLOW1)

# eval graph
if os.path.exists(EVAL_GRAPH):
    slide.shapes.add_picture(EVAL_GRAPH, Inches(0.5), Inches(1.3), Inches(6.5), Inches(3.8))

# Issues
issues = [
    ("問題一：過擬合", "Val Loss 在第 3-4 輪後反彈，模型死記訓練集"),
    ("問題二：重複生成", "部分輸出陷入重複循環（如重複問題本身）"),
    ("問題三：幻覺", "模型輸出與 OpenClaw 無關的錯誤知識"),
    ("問題四：領域泛化不足", "207 筆資料無法覆蓋足夠多的問法變體"),
]
for i, (title, desc) in enumerate(issues):
    y = Inches(5.3 + i * 0.5)
    add_shape(slide, Inches(0.8), y, Inches(0.08), Inches(0.4), YELLOW1, MSO_SHAPE.RECTANGLE)
    add_text_box(slide, Inches(1.1), y, Inches(11), Inches(0.4),
                 f"{title}：{desc}", font_size=15, color=DARK)


# ============ Slide 8: Improvements ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), YELLOW1)
add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), BLUE2, MSO_SHAPE.RECTANGLE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "改進方向", font_size=36, bold=True, color=BLUE2)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), YELLOW1)

improvements = [
    ("增加資料量", "目標 500-1000 筆，涵蓋更多問法變體和邊界情況"),
    ("資料擴增", "使用 LLM 將現有 QA 對改寫為不同表述，提升泛化性"),
    ("降低 LoRA Rank", "r=4, alpha=8，減少過擬合風險"),
    ("Early Stopping", "設置 patience=3，val loss 不改善即停止訓練"),
    ("提高 Dropout", "LoRA dropout 從 0.05 提高到 0.1-0.2"),
    ("QLoRA", "使用 4-bit 量化 + LoRA，可用更大 batch size"),
]
for i, (k, v) in enumerate(improvements):
    x = Inches(0.8) if i < 3 else Inches(7)
    y = Inches(1.8 + (i % 3) * 1.6)
    card = add_shape(slide, x, y, Inches(5.5), Inches(1.3), YELLOW2)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(4.9), Inches(0.4),
                 k, font_size=20, bold=True, color=BLUE2)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.6), Inches(4.9), Inches(0.6),
                 v, font_size=15, color=DARK)


# ============ Slide 9: API Service ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), YELLOW1)
add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), BLUE2, MSO_SHAPE.RECTANGLE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "API 服務 & 實戰應用", font_size=36, bold=True, color=BLUE2)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), YELLOW1)

# Endpoint sections
api_endpoints = [
    ("POST /v1/chat/completions", "OpenAI 兼容對話接口，支援 streaming"),
    ("GET /v1/models", "列出可用模型"),
    ("GET /health", "健康檢查"),
]
for i, (ep, desc) in enumerate(api_endpoints):
    y = Inches(1.6 + i * 0.7)
    add_shape(slide, Inches(0.8), y, Inches(5), Inches(0.5), BLUE3)
    add_text_box(slide, Inches(1), y + Inches(0.05), Inches(4.6), Inches(0.4),
                 ep, font_size=16, bold=True, color=WHITE)
    add_text_box(slide, Inches(6.2), y + Inches(0.05), Inches(6), Inches(0.4),
                 desc, font_size=16, color=DARK)

# Curl example
code_box = add_shape(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.5), RGBColor(0x2D, 0x2D, 0x2D))
add_text_box(slide, Inches(1.2), Inches(4.4), Inches(10.8), Inches(2.3),
             'curl http://localhost:8000/v1/chat/completions \\\n'
             '  -H "Content-Type: application/json" \\\n'
             '  -d \'{\n'
             '    "messages": [{"role": "user", "content": "OpenClaw 是什麼？"}],\n'
             '    "max_tokens": 200, "temperature": 0.7\n'
             '  }\'',
             font_size=14, color=YELLOW1)

add_text_box(slide, Inches(0.8), Inches(6.9), Inches(11), Inches(0.5),
             "部署方式：CUDA_VISIBLE_DEVICES=6 python api_server.py    運行在 0.0.0.0:8000",
             font_size=14, color=BLUE2)


# ============ Slide 10: Summary ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Full yellow top
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(3.2), YELLOW1)
# Full blue bottom
add_shape(slide, Inches(0), Inches(3.2), Inches(13.33), Inches(4.3), BLUE2)

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(1),
             "總結", font_size=40, bold=True, color=DARK, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "從資料收集 → LoRA 微調 → 評估分析 → API 部署，完整 NLP 模型迭代流程",
             font_size=20, color=DARK, align=PP_ALIGN.CENTER)

summary_items = [
    "✓ 成功完成 Qwen3.5-0.8B 在 OpenClaw FAQ 上的 LoRA 微調",
    "✓ Train Loss 從 2.93 降至 0.013，模型有效學習了訓練資料",
    "⚠ Val Loss 第 4 輪後反彈，存在過擬合問題需改進",
    "✓ 提供 OpenAI 兼容 API，可直接整合到 agent 系統",
    "→ 下一步：擴充資料集 + 調整超參數 + 嘗試 QLoRA",
]
for i, item in enumerate(summary_items):
    y = Inches(3.6 + i * 0.65)
    prefix = item[:1]
    color = RGBColor(0x2E, 0xCC, 0x71) if prefix in ("✓",) else (YELLOW1 if prefix == "⚠" else BLUE3)
    add_shape(slide, Inches(1.2), y, Inches(0.06), Inches(0.45), color, MSO_SHAPE.RECTANGLE)
    add_text_box(slide, Inches(1.6), y, Inches(10), Inches(0.5),
                 item, font_size=18, color=WHITE if i > 1 else DARK)


# ============ Save ============
output_path = os.path.join(BASE_DIR, "Qwen3.5_FAQ_微調項目簡報.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
